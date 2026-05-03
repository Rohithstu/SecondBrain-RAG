"""
SecondBrain Core Engine (Cycle 5)
Advanced RAG with Metadata Enrichment & LLM Generation.
- Automated Topic/Keyword/Risk Extraction (Cycle 5)
- LLM-Grounded Answer Generation
- FAISS Vector Management with Metadata Persistence
"""

import os
import re
import json
import hashlib
import numpy as np # type: ignore
import faiss # type: ignore
import torch # type: ignore
import time
import requests # type: ignore
from sentence_transformers import SentenceTransformer, util # type: ignore
from threading import Thread
import google.generativeai as genai # type: ignore
from typing import List, Dict, Set, Any, cast, Tuple, Optional, Union, SupportsIndex

# Format Support Imports
import PyPDF2 # type: ignore
from docx import Document as DocxDocument # type: ignore
from pptx import Presentation # type: ignore
import pandas as pd # type: ignore
from PIL import Image # type: ignore
import pytesseract # type: ignore

# Monitoring
from watchdog.observers import Observer # type: ignore
from watchdog.events import FileSystemEventHandler # type: ignore


# ── LLM PROVIDERS ─────────────────────────────────────────────────────────────

class BaseLLMProvider:
    def generate(self, prompt: str) -> str:
        raise NotImplementedError

    def extract_metadata(self, text: str) -> Dict[str, Any]:
        raise NotImplementedError


class GeminiProvider(BaseLLMProvider):
    def __init__(self, api_key: str):
        genai.configure(api_key=api_key)
        self.model_name = "gemini-2.5-flash" 
        self.model = genai.GenerativeModel(self.model_name)
        print(f"[Gemini] Initialized with model: {self.model_name}")

    def generate(self, prompt: str) -> str:
        for attempt in range(3):
            try:
                response = self.model.generate_content(prompt)
                if hasattr(response, "text"):
                    return str(response.text)
                return "Error: No text in response"
            except Exception as e:
                # Handle 429 Quota Exceeded with backoff
                if "429" in str(e):
                    wait = (attempt + 1) * 3
                    print(f"[Gemini] Quota hit, waiting {wait}s...")
                    time.sleep(wait)
                    continue
                print(f"[Gemini] Generation error: {e}")
                return f"Error generating answer: {str(e)}"
        return "Service temporarily unavailable due to quota limits. Please try again in a few minutes."

    def extract_metadata(self, text: str) -> Dict[str, Any]:
        prompt = f"""
        Analyze the following text and extract metadata in JSON format ONLY:
        Text: {text}
        
        Required JSON Fields:
        - topics: list of 3-5 main topics
        - keywords: list of 5-8 key terms
        - summary: a 2-sentence concise summary
        - risks: any potential risks or conflicts mentioned (if none, return empty list)
        
        JSON:
        """
        try:
            response = self.model.generate_content(prompt)
            raw_text = response.text
            # Simple JSON extraction
            json_match = re.search(r'\{.*\}', raw_text, re.DOTALL)
            if json_match:
                return json.loads(json_match.group(0))
            return {"topics": ["General"], "keywords": [], "summary": "No summary generated.", "risks": []}
        except Exception as e:
            print(f"[Gemini] Metadata error: {e}")
            return {"topics": ["General"], "keywords": [], "summary": "N/A", "risks": []}


class MockLLMProvider(BaseLLMProvider):
    """Fallback for local testing without API keys."""
    def generate(self, prompt: str) -> str:
        return "This is a mock answer based on the provided context. (Add API key for real generation)"

    def extract_metadata(self, text: str) -> Dict[str, Any]:
        return {
            "topics": ["Analysis Pending"],
            "keywords": ["Local", "Knowledge"],
            "summary": "This document has been indexed locally.",
            "risks": []
        }


# ── CORE ENGINE ─────────────────────────────────────────────────────────────

class SecondBrainEngine:
    def __init__(
        self,
        user_id: str,
        base_data_folder: str = "data",
        model_name: str = "all-MiniLM-L6-v2",
        relevance_threshold: float = 0.40,
        shared_model: Optional[SentenceTransformer] = None
    ):
        self.user_id = str(user_id)
        self.user_folder = os.path.join(base_data_folder, self.user_id)
        self.docs_folder = os.path.join(self.user_folder, "docs")
        self.index_file = os.path.join(self.user_folder, "index.faiss")
        self.metadata_file = os.path.join(self.user_folder, "metadata.json")
        self.relevance_threshold = relevance_threshold
        
        os.makedirs(self.docs_folder, exist_ok=True)
        
        # Share the embedding model across user engines to save memory
        if shared_model:
            self.model = shared_model
        else:
            print(f"[Engine-{self.user_id}] Initializing Embedding Model: {model_name}...")
            self.model = SentenceTransformer(model_name)
        
        # LLM Initialization
        api_key = os.getenv("GEMINI_API_KEY")
        if api_key:
            self.llm = GeminiProvider(api_key)
        else:
            self.llm = MockLLMProvider()
        
        self.all_chunks: List[str] = []
        self.chunk_sources: List[str] = []
        self.file_metadata: Dict[str, Dict[str, Any]] = {}
        self.index: Optional[faiss.IndexFlatL2] = None
        
        self._load_from_disk()
        self.refresh_index()
        self._patch_missing_metadata()

    def _patch_missing_metadata(self) -> None:
        """Heuristic: if topics are missing, the file needs enrichment."""
        to_patch = [rel for rel, meta in self.file_metadata.items() if "topics" not in meta]
        if not to_patch: return
        
        print(f"[Engine-{self.user_id}] Patching metadata for {len(to_patch)} files...")
        for rel in to_patch:
            full = os.path.join(self.docs_folder, rel)
            ext = rel.lower().split(".")[-1] if "." in rel else ""
            content = ""
            try:
                if ext == "txt":
                    with open(full, "r", encoding="utf-8", errors="ignore") as f: content = f.read()
                else: content = self._load_formats(full, ext)
                
                if content.strip():
                    cleaned = self._advanced_clean(content)
                    doc_meta = self._extract_file_metadata(cleaned, rel)
                    self.file_metadata[rel].update(doc_meta)
            except Exception as e:
                print(f"  [Patch] Failed for {rel}: {e}")
        
        self._save_to_disk()

    # ── CYCLE 5: ADVANCED METADATA ENRICHMENT ──────────────────────────────
    def _extract_file_metadata(self, text: str, rel_path: str) -> Dict[str, Any]:
        """Extracts topics, keywords, and summary using the LLM."""
        print(f"  [Enrichment-{self.user_id}] Analyzing {rel_path}...")
        snippet = text[:4000]
        meta = cast(Dict[str, Any], self.llm.extract_metadata(snippet))
        meta["last_indexed"] = time.time()
        return meta

    # ── CYCLE 3/4 REFINEMENTS ─────────────────────────────────────────────
    def _advanced_clean(self, text: str) -> str:
        text = re.sub(r'(UNIT|Unit|Slide|SLIDE|PAGE|Page)\s+\d+', '', text)
        text = re.sub(r'Author:.*|Bharath\s+Yannam', '', text, flags=re.IGNORECASE)
        text = re.sub(r'(\w+)-\s*\n\s*(\w+)', r'\1\2', text)
        return " ".join(text.split())

    def _chunk_text(self, text: str) -> List[str]:
        """Splits text into overlapping chunks for better retrieval."""
        raw_sentences: List[Any] = re.split(r'(?<=[.!?])\s+', text)
        sentences: List[str] = [str(s).strip() for s in raw_sentences if s and len(str(s).strip()) > 5]
        chunks: List[str] = []
        for i in range(0, len(sentences), 2):
            limit = i + 3
            chunk_sents = [sentences[j] for j in range(i, min(limit, len(sentences)))]
            if not chunk_sents: continue
            chunks.append(" ".join(chunk_sents))
        return chunks

    # ── PERSISTENCE ──────────────────────────────────────────────────────────
    def _get_file_hash(self, filepath: str) -> str:
        hasher = hashlib.md5()
        try:
            with open(filepath, 'rb') as f: hasher.update(f.read())
            return hasher.hexdigest()
        except: return ""

    def _load_from_disk(self) -> None:
        if os.path.exists(self.index_file) and os.path.exists(self.metadata_file):
            try:
                self.index = faiss.read_index(self.index_file)
                with open(self.metadata_file, "r", encoding="utf-8") as f:
                    data = json.load(f)
                    self.all_chunks = data.get("all_chunks", [])
                    self.chunk_sources = data.get("chunk_sources", [])
                    self.file_metadata = data.get("file_metadata", {})
            except Exception as e:
                print(f"[Engine-{self.user_id}] Load error: {e}")
                self.index = None

    def _save_to_disk(self) -> None:
        if self.index is not None:
            faiss.write_index(self.index, self.index_file)
        with open(self.metadata_file, "w", encoding="utf-8") as f:
            json.dump({
                "all_chunks": self.all_chunks, 
                "chunk_sources": self.chunk_sources, 
                "file_metadata": self.file_metadata
            }, f, indent=2)

    def _load_formats(self, path: str, ext: str) -> str:
        try:
            if ext == "pdf":
                with open(path, "rb") as f:
                    return " ".join(p.extract_text() or "" for p in PyPDF2.PdfReader(f).pages)
            if ext == "docx":
                doc = DocxDocument(path)
                return "\n".join(p.text for p in doc.paragraphs)
            if ext == "pptx":
                t_builder = []
                prs = Presentation(path)
                for s in prs.slides:
                    for sh in s.shapes:
                        t = getattr(sh, "text", "")
                        if t: t_builder.append(str(t))
                return " ".join(t_builder)
            if ext == "xlsx":
                df_d = pd.read_excel(path, sheet_name=None)
                return " ".join(str(df.to_string()) for df in df_d.values())
            if ext in ["png", "jpg", "jpeg"]:
                img = Image.open(path)
                text = str(pytesseract.image_to_string(img, config='--psm 11'))
                return text
        except Exception as e:
            print(f"  [Format Error] Failed to load {ext}: {e}")
        return ""

    def refresh_index(self) -> None:
        print(f"[Engine-{self.user_id}] Refreshing Index...")
        c_state = {}
        for r, ds, fs in os.walk(self.docs_folder):
            for f in fs:
                p = os.path.join(r, f)
                rel = os.path.relpath(p, self.docs_folder)
                c_state[rel] = {"hash": self._get_file_hash(p)}

        changed = set(rel for rel in self.file_metadata if rel not in c_state or self.file_metadata[rel].get("hash") != c_state[rel]["hash"])
        if changed:
            pac: List[str] = [str(x) for x in self.all_chunks]
            pcs: List[str] = [str(x) for x in self.chunk_sources]
            self.all_chunks, self.chunk_sources = [], []
            for c, s in zip(pac, pcs):
                if str(s) not in changed:
                    self.all_chunks.append(str(c))
                    self.chunk_sources.append(str(s))
            for s in changed: self.file_metadata.pop(str(s), None)
            
            if self.all_chunks:
                embs = np.array(self.model.encode(self.all_chunks, show_progress_bar=False), dtype=np.float32)
                self.index = cast(faiss.IndexFlatL2, faiss.IndexFlatL2(embs.shape[1]))
                getattr(self.index, "add")(embs)
            else:
                self.index = None

        modified = False
        for rel in c_state:
            if rel not in self.file_metadata:
                full = os.path.join(self.docs_folder, rel)
                ext = rel.lower().split(".")[-1] if "." in rel else ""
                content = ""
                if ext == "txt":
                    with open(full, "r", encoding="utf-8", errors="ignore") as f: content = f.read()
                else: content = self._load_formats(full, ext)
                
                cleaned = self._advanced_clean(content)
                doc_meta: Dict[str, Any] = {}
                if cleaned.strip():
                    doc_meta = self._extract_file_metadata(cleaned, rel)
                else:
                    doc_meta = {"topics": ["Unreadable"], "keywords": [], "summary": "N/A", "risks": []}
                
                doc_meta["hash"] = c_state[rel]["hash"]
                self.file_metadata[rel] = doc_meta
                modified = True
                
                if cleaned.strip():
                    new_chunks = self._chunk_text(cleaned)
                    if new_chunks:
                        embs = np.array(self.model.encode(new_chunks, show_progress_bar=False), dtype=np.float32)
                        if self.index is None: 
                            self.index = cast(faiss.IndexFlatL2, faiss.IndexFlatL2(embs.shape[1]))
                        getattr(self.index, "add")(embs)
                        self.all_chunks.extend(new_chunks)
                        self.chunk_sources.extend([str(rel)] * len(new_chunks))
        
        if modified or changed:
            self._save_to_disk()

    # ── SEARCH & GENERATION ───────────────────────────────────────────────
    def search(self, query: str, top_k: int = 6, offline: bool = False) -> Dict[str, Any]:
        if offline: top_k = 10
        if self.index is None or not self.all_chunks:
            return {"answer": "I don't have any knowledge yet. Please upload documents.", "confidence": 0}

        q_emb = np.array(self.model.encode([query]), dtype=np.float32)
        distances, indices = getattr(self.index, "search")(q_emb, top_k)
        
        context_chunks = []
        sources = set()
        dist_row = cast(List[float], distances[0])
        idx_row = cast(List[int], indices[0])
        
        for d, i in zip(dist_row, idx_row):
            idx = int(i)
            if idx == -1: continue
            score = 1.0 / (1.0 + float(d))
            if score < self.relevance_threshold: continue
            
            chunk_text = str(self.all_chunks[idx])
            source = str(self.chunk_sources[idx])
            context_chunks.append(f"[Source: {source}] {chunk_text}")
            sources.add(source)

        if not context_chunks:
            return {"answer": "No relevant content found in your documents.", "confidence": 0}

        if offline:
            return self._offline_synthesize(query, context_chunks, list(sources))

        context_block = "\n---\n".join(context_chunks)
        prompt = f"""
        STRICT GROUNDING INSTRUCTIONS:
        Use ONLY the provided context blocks to answer the user query.
        Context: {context_block}
        Query: {query}
        Rules: Zero hallucination, Zero external knowledge, Markdown formatting, Inline citations [filename].
        """
        answer = self.llm.generate(prompt)
        
        featured_topics = []
        for s in sources:
            if s in self.file_metadata:
                featured_topics.extend(self.file_metadata[s].get("topics", []))
        
        return {
            "answer": str(answer),
            "sources": list(sources),
            "topics": list(set(featured_topics))[:5],
            "confidence": 0.9,
            "status": "LLM Generated"
        }

    def _offline_synthesize(self, query: str, context_chunks: List[str], sources: List[str]) -> Dict[str, Any]:
        # Simple synthesis for offline
        best_chunk = context_chunks[0]
        return {
            "answer": f"### Offline Answer\n{best_chunk}",
            "sources": sources,
            "confidence": 0.7,
            "status": "Offline Match"
        }

# ── ENGINE MANAGER ────────────────────────────────────────────────────────
class SecondBrainManager:
    """Manages multiple SecondBrainEngine instances (one per user)."""
    def __init__(self, base_data_folder: str = "data"):
        self.base_data_folder = base_data_folder
        self.shared_model = SentenceTransformer("all-MiniLM-L6-v2")
        self.user_engines: Dict[str, SecondBrainEngine] = {}

    def get_engine(self, user_id: Union[str, int]) -> SecondBrainEngine:
        uid = str(user_id)
        if uid not in self.user_engines:
            self.user_engines[uid] = SecondBrainEngine(
                user_id=uid,
                base_data_folder=self.base_data_folder,
                shared_model=self.shared_model
            )
        return self.user_engines[uid]

# ── MONITORING ─────────────────────────────────────────────────────────────
# ── MONITORING ─────────────────────────────────────────────────────────────
class DataMonitorHandler(FileSystemEventHandler):
    def __init__(self, manager: SecondBrainManager): 
        self.manager = manager
    def on_any_event(self, event):
        if not event.is_directory and not event.src_path.endswith((".json", ".faiss")):
            # Path structure: data/<user_id>/docs/file.ext
            norm_path = os.path.normpath(event.src_path)
            parts = norm_path.split(os.sep)
            try:
                # Find the directory after 'data'
                if self.manager.base_data_folder in parts:
                    idx = parts.index(self.manager.base_data_folder)
                    if len(parts) > idx + 1:
                        user_id = parts[idx + 1]
                        time.sleep(1)
                        self.manager.get_engine(user_id).refresh_index()
            except Exception: pass

def start_monitoring(manager: SecondBrainManager):
    handler = DataMonitorHandler(manager)
    ob = Observer()
    ob.schedule(handler, manager.base_data_folder, recursive=True)
    ob.start(); return ob
